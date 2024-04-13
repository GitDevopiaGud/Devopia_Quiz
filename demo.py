from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import requests
import pdfplumber
from docx import Document
from pptx import Presentation

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

UPLOAD_FOLDER = 'uploads/'
ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.pptx', '.doc', '.ppt'}
# Define the target API endpoint
TARGET_API_URL = 'https://api.example.com/endpoint'

# Ensure the upload folder exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return os.path.splitext(filename)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/files/upload', methods=['POST'])
def upload_file():
    file = request.files.get('files')
    if not file or not allowed_file(file.filename):
        return jsonify({'error': 'Unsupported file type or no file uploaded'}), 400
    
    # Save the uploaded file to the upload folder
    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)
    
    file_extension = os.path.splitext(file.filename)[1].lower()
    text_content = ""

    try:
        if file_extension == '.pdf':
            # Use pdfplumber to parse PDF file
            with pdfplumber.open(file_path) as pdf:
                text_content = "\n".join([page.extract_text() for page in pdf.pages])
        
        elif file_extension == '.docx' or file_extension == '.doc':
            # Use python-docx to parse DOCX file
            doc = Document(file_path)
            text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        elif file_extension == '.txt':
            # Read and extract text from TXT file
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
        
        elif file_extension == '.pptx' or file_extension == '.ppt':
            # Use python-pptx to parse PPT file
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text_content += shape.text + "\n"
        
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        return jsonify({'text_content': text_content}), 200
        # Send text content to another API
        payload = {'text': text_content}
        response = requests.post(TARGET_API_URL, json=payload)
                
        # Remove the uploaded file
        os.remove(file_path)
        
        # Return the response from the API or customize your response as needed
        return jsonify({'message': 'File processed successfully', 'api_response': response.json()}), 200
    
    except Exception as error:
        print(error)
        return jsonify({'error': 'An error occurred'}), 500

if __name__ == '__main__':
    app.run(debug=True)
