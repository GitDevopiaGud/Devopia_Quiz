import re
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from collections import Counter
import nltk
import pdfplumber
from docx import Document
from pptx import Presentation

nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')

def remove_duplicates(text):
    sentences = nltk.sent_tokenize(text)
    sentence_counts = Counter(sentences)
    unique_sentences = list(sentence_counts.keys())
    return ' '.join(unique_sentences)

def remove_stopwords(text):
    stop_words = set(stopwords.words('english'))
    word_tokens = word_tokenize(text)
    filtered_text = [word for word in word_tokens if word.lower() not in stop_words]
    return ' '.join(filtered_text)

def compress_sentences(text):
    sentences = nltk.sent_tokenize(text)
    lemmatizer = WordNetLemmatizer()
    lemmatized_sentences = [' '.join(lemmatizer.lemmatize(word) for word in nltk.word_tokenize(sentence)) for sentence in sentences]
    return ' '.join(lemmatized_sentences)

def extract_text_from_ppt(ppt_file):
    text = ""
    presentation = Presentation(ppt_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(pdf_file):
    text = ""
    with pdfplumber.open(pdf_file) as pdf:
        text = "\n".join([page.extract_text() for page in pdf.pages])
    return text

def extract_text_from_docx(docx_file):
    text = ""
    doc = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_txt(txt_file):
    text = ""
    with open(txt_file, 'r', encoding='utf-8') as f:
        text = f.read()
    return text

def clean_text(text):
    # Remove unwanted characters such as escape sequences and unwanted symbols
    cleaned_text = re.sub(r'[\x80-\xFF]+', '', text)  # Remove Unicode characters in the range 0x80-0xFF
    cleaned_text = re.sub(r'\\[nrt]', '', cleaned_text)  # Remove escaped newline, tab, and return characters
    cleaned_text = re.sub(r'[^\x00-\x7F]+', ' ', cleaned_text)  # Replace non-ASCII characters with space
    cleaned_text = re.sub(r' +', ' ', cleaned_text)  # Replace multiple spaces with a single space
    cleaned_text = cleaned_text.strip()  # Remove leading and trailing whitespace
    return cleaned_text

def extract_text(file_path):
    if file_path.endswith('.pptx'):
        return extract_text_from_ppt(file_path)
    elif file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    else:
        raise ValueError("Unsupported file format")

file_path = "/Users/chiragsharma/Downloads/SDC_4.2_IO_System (1).pdf"
file_text = extract_text(file_path)

file_text = remove_duplicates(file_text)
file_text = remove_stopwords(file_text)
file_text = compress_sentences(file_text)
def decode_with_space_replace(bytes_obj):
    bytes_obj.decode('utf-8', errors='replace')
    return bytes_obj

text_content = decode_with_space_replace(file_text.encode(errors='ignore'))
# Remove the unnecessary encoding step
# file_text is already a string type
# We can directly pass it to clean_text
cleaned_text = clean_text(file_text)

# Print the cleaned text
print(cleaned_text)
