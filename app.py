from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import json
import pymongo
import requests
import pdfplumber
from docx import Document
from pptx import Presentation
import re
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from collections import Counter
import nltk

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes
nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')

# Configure MongoDB connection
client = pymongo.MongoClient("mongodb+srv://sharmachirag393:JlOSu0BFSiJ7FnRF@devopia.qwmpayp.mongodb.net/?retryWrites=true&w=majority&appName=Devopia")
db = client["test"]
user_collection = db["users"]
quiz_collection = db["quiz"]
userschema_collection = db["userschema"]

UPLOAD_FOLDER = 'uploads/'
ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.pptx', '.doc', '.ppt'}
# Define the target API endpoint
TARGET_API_URL = 'https://expressllm.onrender.com/generate'

# Ensure the upload folder exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return os.path.splitext(filename)[1].lower() in ALLOWED_EXTENSIONS

# def remove_duplicates(text):
#     sentences = nltk.sent_tokenize(text)
#     sentence_counts = Counter(sentences)
#     unique_sentences = list(sentence_counts.keys())
#     return ' '.join(unique_sentences)

# def remove_stopwords(text):
#     stop_words = set(stopwords.words('english'))
#     word_tokens = word_tokenize(text)
#     filtered_text = [word for word in word_tokens if word.lower() not in stop_words]
#     return ' '.join(filtered_text)

# def compress_sentences(text):
#     sentences = nltk.sent_tokenize(text)
#     lemmatizer = WordNetLemmatizer()
#     lemmatized_sentences = [' '.join(lemmatizer.lemmatize(word) for word in nltk.word_tokenize(sentence)) for sentence in sentences]
#     return ' '.join(lemmatized_sentences)

def remove_duplicates(text):
    # Tokenize text into sentences
    sentences = nltk.sent_tokenize(text)
    
    # Use a set to track unique sentences and preserve order
    unique_sentences = []
    seen_sentences = set()
    
    for sentence in sentences:
        if sentence not in seen_sentences:
            unique_sentences.append(sentence)
            seen_sentences.add(sentence)
    
    # Join unique sentences back into a single text
    return ' '.join(unique_sentences)

def remove_stopwords(text):
    # Get the set of stopwords
    stop_words = set(stopwords.words('english'))
    
    # Tokenize the text into words
    words = word_tokenize(text)
    
    # Remove stopwords and reconstruct the text
    filtered_words = [word for word in words if word.lower() not in stop_words]
    filtered_text = ' '.join(filtered_words)
    
    return filtered_text

def compress_sentences(text):
    lemmatizer = WordNetLemmatizer()
    # Split text into sentences
    sentences = nltk.sent_tokenize(text)
    compressed_sentences = []
    
    for sentence in sentences:
        # Tokenize sentence and lemmatize each word
        words = nltk.word_tokenize(sentence)
        lemmatized_words = [lemmatizer.lemmatize(word) for word in words]
        # Join lemmatized words to reconstruct sentence
        compressed_sentence = ' '.join(lemmatized_words)
        compressed_sentences.append(compressed_sentence)
    
    # Join all compressed sentences to form the final text
    return ' '.join(compressed_sentences)


def extract_text_from_ppt(ppt_file):
    text = ""
    presentation = Presentation(ppt_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(pdf_file):
    text = ""
    with pdfplumber.open(pdf_file) as pdf:
        text = "\n".join([page.extract_text() for page in pdf.pages])
    return text

def extract_text_from_docx(docx_file):
    text = ""
    doc = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def extract_text_from_txt(txt_file):
    text = ""
    with open(txt_file, 'r', encoding='utf-8') as f:
        text = f.read()
    return text

def clean_text(text):
    # Remove unwanted characters such as escape sequences and unwanted symbols
    cleaned_text = re.sub(r'[\x80-\xFF]+', '', text)  # Remove Unicode characters in the range 0x80-0xFF
    cleaned_text = re.sub(r'\\[nrt]', '', cleaned_text)  # Remove escaped newline, tab, and return characters
    cleaned_text = re.sub(r'[^\x00-\x7F]+', ' ', cleaned_text)  # Replace non-ASCII characters with space
    cleaned_text = re.sub(r' +', ' ', cleaned_text)  # Replace multiple spaces with a single space
    cleaned_text = cleaned_text.strip()  # Remove leading and trailing whitespace
    return cleaned_text

def extract_text(file_path):
    if file_path.endswith('.pptx'):
        return extract_text_from_ppt(file_path)
    elif file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    else:
        raise ValueError("Unsupported file format")

@app.route('/')
def home():
    return 'Welcome to the File Upload API'

@app.route('/files/upload', methods=['POST'])
def upload_file():
    file = request.files.get('files')
    email = request.form.get('email')
    if not file or not allowed_file(file.filename):
        return jsonify({'error': 'Unsupported file type or no file uploaded'}), 400
    
    # Save the uploaded file to the upload folder
    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)
    file_text = ""
    try:
        file_text = extract_text(file_path)

        # Process the file text
        file_text = remove_duplicates(file_text)
        file_text = remove_stopwords(file_text)
        file_text = compress_sentences(file_text)

        # Clean the text content
        cleaned_text = clean_text(file_text)
        # Send text content to another API
        payload = {'text': cleaned_text}
        response = requests.post(TARGET_API_URL, json=payload)
        
        api_response = response.json()
        mcq_data = json.loads(api_response['mcq'].replace('\\n', '\n'))
        title = api_response['title']
        questions = [item['question'] for item in mcq_data]
        options = [[option for option in item['options']] for item in mcq_data]
        correct_answer = [item['correctAnswer'] for item in mcq_data]
        
        # Remove the uploaded file
        os.remove(file_path)
        
        return jsonify({'email': email, 'title': title, 'questions': questions, 'options': options, 'correct_answer': correct_answer}), 200
    
    except Exception as error:
        print(error)
        # Return a JSON response indicating an error
        return jsonify({'error': 'An error occurred during file processing'}), 500
    
@app.route('/files/database', methods=['POST'])
def upload_to_database():
    # Get JSON data from the request body
    data = request.get_json()
    
    if not data:
        return jsonify({'error': 'Invalid or missing JSON data'}), 400
    
    # Retrieve individual fields from the JSON data
    email = data.get('email')
    title = data.get('title')
    questions = data.get('questions')
    options = data.get('options')
    correct_answer = data.get('correct_answer')
    
    # Validate required fields
    if not all([email, title, questions, options, correct_answer]):
        return jsonify({'error': 'Missing required fields'}), 400
    
    # Fetch the user from the User collection using email
    user = user_collection.find_one({'email': email})
    
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    grade = user.get('current_class')
    board = user.get('board')
    
    # Create a new Quiz document
    new_quiz = {
        'email': email,
        'topic_name': title,
        'grade': grade,
        'board': board,
        'questions': questions,
        'options': options,
        'correct_answer': correct_answer
    }
    
    # Insert the new quiz document into the quiz_collection
    quiz_collection.insert_one(new_quiz)
    
    return jsonify({'message': 'Quiz data saved successfully'}), 200

@app.route('/files/userschema', methods=['POST'])
def upload_to_user():
    # Get JSON data from the request body
    data = request.get_json()
    
    if not data:
        return jsonify({'error': 'Invalid or missing JSON data'}), 400
    
    # Retrieve individual fields from the JSON data
    email = data.get('email')
    quiz_id = data.get('quiz_id')
    options = data.get('options')
    correct_answer = data.get('correct_answer')
    
    # Validate required fields
    if not all([email, quiz_id, options, correct_answer]):
        return jsonify({'error': 'Missing required fields'}), 400
    
    # Fetch the user from the User collection using email
    user = user_collection.find_one({'email': email})
    
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    unique = user.get('_id')
    
    # Create a new Quiz document
    new_quiz = {
        'email_id': unique,
        'quiz_id': quiz_id,
        'options': options,
        'correct_answer': correct_answer
    }
    
    # Insert the new quiz document into the userschema_collection
    userschema_collection.insert_one(new_quiz)
    
    return jsonify({'message': 'User data saved successfully'}), 200

if __name__ == '__main__':
    app.run(debug=True)
